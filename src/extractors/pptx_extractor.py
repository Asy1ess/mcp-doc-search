"""PowerPoint PPTX text extraction."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.shapes.base import BaseShape

from src.extractors._utils import make_document


def _shape_text(shape: BaseShape) -> str:
    if not hasattr(shape, "text"):
        return ""
    text = shape.text.strip()  # type: ignore[union-attr]
    return text


def extract_pptx(path: Path) -> ExtractedDocument:
    presentation = Presentation(path)
    slides: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        parts = [_shape_text(shape) for shape in slide.shapes]
        parts = [p for p in parts if p]
        if parts:
            slides.append(f"## Slide {index}\n" + "\n".join(parts))

    return make_document(
        path,
        "\n\n".join(slides),
        format="pptx",
        slide_count=str(len(presentation.slides)),
    )
